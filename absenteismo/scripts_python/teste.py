

from datetime import date
import sys
from pathlib import Path
from dateutil.relativedelta import relativedelta
import smtplib
import time
import os
from dotenv import load_dotenv

# Path(__file__).resolve()        = caminho completo até grafico_linha_absenteismo_por_mes.py
# .parent                         = pasta 'geral/'
# .parent.parent                  = pasta raiz 'Scripts_python/'
DIRETORIO_RAIZ = Path(__file__).resolve().parent

# Adiciona a raiz do projeto ao caminho de busca do Python
if str(DIRETORIO_RAIZ) not in sys.path:
    sys.path.append(str(DIRETORIO_RAIZ))


from pptx import Presentation
from pptx.util import Inches

prs = Presentation("/home/suporte/scripts/rh/absenteismo/assets/template_apresentacao_rh.pptx")

slide = prs.slides[0] # Altere para o índice do slide que deseja verificar
for shape in slide.placeholders:
    print(f"Index: {shape.placeholder_format.idx} | Nome: {shape.name} | Texto Atual: '{shape.text}'")

