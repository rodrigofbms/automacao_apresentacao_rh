import logging
from pptx import Presentation
from pptx.util import Pt
from PIL import Image as PILImage
from pathlib import Path

# ---------------------------------------------------------
# 4. IMPORTS DA CONFIGURAÇÕES
# ---------------------------------------------------------
from config.settings import MES_ANO_UNDERLINE, MES_ANO_FILTRO, BASE_DIR
from config.pptx_mapping import obter_mapa_titulos, obter_mapa_imagens, obter_mapa_texto_ai

# ---------------------------------------------------------
# BUSCAR TODAS AS IMAGENS DO MÊS NAS SUBPASTAS
# ---------------------------------------------------------
def coletar_imagens_graficos() -> list[Path]:
    """
    Procura recursivamente por arquivos .png que contenham a tag do mês
    (ex: 2026_06 ou 2026-06) em qualquer subpasta de 'Scripts_python/'.
    """
    imagens_encontradas = []

    # Faz uma busca recursiva em todas as subpastas
    for arquivo in BASE_DIR.rglob("*.png"):
        nome_arquivo = arquivo.name
        
        # Filtra imagens do mês de referência (seja com _ ou -)
        if MES_ANO_UNDERLINE in nome_arquivo or MES_ANO_FILTRO in nome_arquivo:
            imagens_encontradas.append(arquivo)

    # Ordena alfabeticamente pelo nome do arquivo/pasta para manter um padrão visual
    imagens_encontradas.sort()
    return imagens_encontradas


# ---------------------------------------------------------
# GERAR APRESENTAÇÃO POWERPOINT (.PPTX)
# ---------------------------------------------------------
def gerar_apresentacao_pptx(lista_imagens, texto_ia: str) -> Path | None:
    """
    Carrega o modelo do PowerPoint e insere os gráficos gerados
    em slides específicos com base no nome do arquivo.
    """
    caminho_template = BASE_DIR.parent / "assets" / "template_apresentacao_rh.pptx"
    pasta_saida = BASE_DIR / "apresentacoes"
    pasta_saida.mkdir(exist_ok=True)
    
    caminho_saida = pasta_saida / f"apresentacao_rh_{MES_ANO_UNDERLINE}.pptx"

    if not caminho_template.exists():
        logging.warning(f"⚠️ Template não encontrado em {caminho_template}. Criando apresentação do zero...")
        return None
    else:
        logging.info(f"📄 Carregando modelo PowerPoint: {caminho_template}")
        presentation = Presentation(caminho_template)

        _inserir_titulos(presentation)
        _inserir_imagens(presentation, lista_imagens)

        if texto_ia:
            _inserir_texto_ia(presentation, texto_ia)

        presentation.save(caminho_saida)
        logging.info(f"✅ Apresentação salva com sucesso em: {caminho_saida}")
        return caminho_saida



# ---------------------------------------------------------
# INSERIR TÍTULOS NA APRESENTAÇÃO DO POWERPOINT (.PPTX)
# ---------------------------------------------------------
def _inserir_titulos(presentation: Presentation):

    mapa_titulos = obter_mapa_titulos()

    # ---------------------------------------------------------
    # ATUALIZAR TÍTULOS DINÂMICOS
    # ---------------------------------------------------------
    for slide_idx, config_titulo in mapa_titulos.items():

        if slide_idx < len(presentation.slides):

            slide = presentation.slides[slide_idx]
            ph_idx = config_titulo["placeholder_index"]
            texto = config_titulo["texto"]

            if config_titulo.get("uppercase", True):
                text = texto.upper()

            try:
                placeholder = slide.placeholders[ph_idx]

                if placeholder.has_text_frame:
                    # Atualiza o texto mantendo a formatação original do template
                    tf = placeholder.text_frame
                    tf.text = texto

                    # Caixas de texto no PowerPoint são divididas em parágrafos. Como títulos geralmente têm apenas 1 linha ou bloco, 
                    # pegamos o primeiro parágrafo (índice 0) para ajustar as propriedades que pertencem ao parágrafo como um todo (como o alinhamento)
                    paragrafo = tf.paragraphs[0]

                    # ----- ALINHAMENTO DO PARÁGRAFO -----
                    if "alinhamento" in config_titulo:
                        # o PP_ALIGN contém as definições e constantes padrões de alinhamento do Power Point, ao invés de passar "center"
                        # como texto, a biblioteca exige a constante formal como PP_ALIGN.CENTER para saber como posicionar o texto na caixa.
                        paragrafo.alignment = config_titulo["alinhamento"]

                    # ----- ESTILIZAÇÃO DO TEXTO -----
                    # No PowerPoint, um parágrafo é dividido em um ou mais Runs (fragmentos de texto que compartilham exatamente a mesma formatação visual).
                    #  Ao pegar o runs[0], acessamos a camada onde a fonte, a cor, o peso (negrito) e o tamanho em pontos realmente residem.
                    if paragrafo.runs:

                        run = paragrafo.runs[0]
                        if "fonte_nome" in config_titulo:
                            run.font.name = config_titulo["fonte_nome"]
                        if "tamanho_pt" in config_titulo:
                            # O power point utiliza um unidade interna de medida. A função Pt() converte o número da fonte padrão (Points)
                            # para o formato que a biblioteca entende, por exemplo: Pt(24) converte para 24 pt.
                            run.font.size = Pt(config_titulo["tamanho_pt"])
                        if "bold" in config_titulo:
                            run.font.bold = config_titulo["bold"]
                        if "italic" in config_titulo:
                            run.font.italic = config_titulo["italic"]
                        if "cor_rgb" in config_titulo:
                            run.font.color.rgb = config_titulo["cor_rgb"]
                
                    logging.info(f"✏️ Título atualizado no Slide {slide_idx + 1}: '{texto}'")

                else:
                    logging.error(f"❌ O placeholder idx {ph_idx} no Slide {slide_idx + 1} não é uma caixa de texto.")

            except KeyError:
                logging.error(f"❌ Placeholder de título idx {ph_idx} não encontrado no Slide {slide_idx + 1}.")

            except Exception as e:
                logging.error(f"❌ Erro ao atualizar título no Slide {slide_idx + 1}: {e}")

        else:
            logging.error(f"❌ Slide {slide_idx + 1} fora do alcance do template (Total: {len(presentation.slides)} slides).")
            break





# ---------------------------------------------------------
# INSERIR IMAGENS NA APRESENTAÇÃO DO POWERPOINT (.PPTX)
# ---------------------------------------------------------
def _inserir_imagens(presentation: Presentation, lista_imagens: list[Path]):

    mapa_imagens = obter_mapa_imagens()

    for img_path in lista_imagens:
        nome_img = img_path.name.lower()
        
        for chave, config in mapa_imagens.items():

            if chave in nome_img:
                slide_idx = config["slide_index"]
                ph_idx = config["placeholder_index"]

                if slide_idx < len(presentation.slides):
                    slide = presentation.slides[slide_idx]
                    try:
                        placeholder = slide.placeholders[ph_idx]

                        # 2. Salva as posições e o tamanho configurado no slide
                        ph_left = placeholder.left
                        ph_top = placeholder.top
                        ph_width = placeholder.width
                        ph_height = placeholder.height
                        ph_ratio = ph_width / ph_height

                        # 3. Descobre o tamanho da imagem gerada pelo Matplotlib/Seaborn
                        with PILImage.open(img_path) as img:
                            img_width, img_height = img.size
                            img_ratio = img_width / img_height

                        # 4. Remove o placeholder antigo para não sobrepor
                        sp = placeholder._element
                        sp.getparent().remove(sp)

                        # 5. Adiciona a imagem no exato local mantendo o preenchimento proporcional (sem esticar)
                        picture = slide.shapes.add_picture(
                            str(img_path), 
                            left=ph_left, 
                            top=ph_top, 
                            width=ph_width, 
                            height=ph_height
                        )
                        """
                        # 2. Ajusta o Crop para preencher o container sem esticar
                        if img_ratio > ph_ratio:
                            # Imagem é mais larga que o placeholder -> corta as laterais
                            new_width = img_height * ph_ratio
                            offset = (img_width - new_width) / 2
                            picture.crop_left = offset / img_width
                            picture.crop_right = offset / img_width
                            picture.crop_top = 0
                            picture.crop_bottom = 0
                        else:
                            # Imagem é mais alta que o placeholder -> corta topo e base
                            new_height = img_width / ph_ratio
                            offset = (img_height - new_height) / 2
                            picture.crop_top = offset / img_height
                            picture.crop_bottom = offset / img_height
                            picture.crop_left = 0
                            picture.crop_right = 0    
                        """
                        logging.info(f"📈 Gráfico '{img_path.name}' inserido no Slide {slide_idx + 1}")

                    except KeyError:
                        logging.error(f"❌ Placeholder idx {ph_idx} não existe no Slide {slide_idx + 1}.")
                    except AttributeError:
                        logging.error(f"❌ O Placeholder idx {ph_idx} no Slide {slide_idx + 1} não aceita inserção de imagem diretamente.")
                    except Exception as e:
                        logging.error(f"❌ Erro ao inserir {img_path.name} no Slide {slide_idx + 1}: {e}")
                else:
                    logging.error(f"❌ Slide {slide_idx + 1} fora do alcance do template (Total: {len(presentation.slides)} slides).")
                break


def _inserir_texto_ia(presentation: Presentation, texto_analise: str):
    """Injeta o parecer gerado pela IA no slide de resumo executivo."""
    config = obter_mapa_texto_ai()
    slide_idx = config["slide_index"]
    ph_idx = config["placeholder_index"]

    if slide_idx >= len(presentation.slides):
        logging.warning(f"⚠️ Slide {slide_idx + 1} para texto da IA não existe na apresentação.")
        return
    else:
        slide = presentation.slides[slide_idx]
        
        try:
            placeholder = slide.placeholders[ph_idx]
            if placeholder.has_text_frame:
                tf = placeholder.text_frame
                tf.text = texto_analise  # Injeta o texto retornado da OpenAI

                # Formatação do texto
                for paragrafo in tf.paragraphs:
                    paragrafo.font.name = config["fonte_nome"]
                    paragrafo.font.size = Pt(config["tamanho_pt"])
                    paragrafo.font.color.rgb = config["cor_rgb"]

                logging.info(f"📝 Análise de IA injetada com sucesso no Slide {slide_idx + 1}.")
        except Exception as e:
            logging.error(f"❌ Erro ao injetar texto da IA no Slide {slide_idx + 1}: {e}")